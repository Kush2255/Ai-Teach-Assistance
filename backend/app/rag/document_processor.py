import os
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
from typing import List, Dict, Any

class DocumentProcessor:
    """Extracts structured text from PDF, DOCX, PPTX, and TXT files retaining page/slide metadata."""

    @staticmethod
    def extract_text(file_path: str) -> List[Dict[str, Any]]:
        ext = os.path.splitext(file_path)[1].lower()
        pages_content = []

        if ext == ".pdf":
            try:
                reader = pypdf.PdfReader(file_path)
                for idx, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages_content.append({
                            "page": idx + 1,
                            "section": f"Page {idx + 1}",
                            "text": text
                        })
            except Exception as e:
                pages_content.append({"page": 1, "section": "Document", "text": f"PDF parse error: {e}"})

        elif ext == ".docx":
            try:
                doc = DocxDocument(file_path)
                current_section = "Chapter 1"
                current_text = []
                page_num = 1
                for para in doc.paragraphs:
                    if para.style.name.startswith("Heading"):
                        if current_text:
                            pages_content.append({
                                "page": page_num,
                                "section": current_section,
                                "text": "\n".join(current_text)
                            })
                            current_text = []
                            page_num += 1
                        current_section = para.text
                    elif para.text.strip():
                        current_text.append(para.text)
                if current_text:
                    pages_content.append({
                        "page": page_num,
                        "section": current_section,
                        "text": "\n".join(current_text)
                    })
            except Exception as e:
                pages_content.append({"page": 1, "section": "Document", "text": f"DOCX parse error: {e}"})

        elif ext == ".pptx":
            try:
                prs = Presentation(file_path)
                for idx, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text)
                    if slide_text:
                        pages_content.append({
                            "page": idx + 1,
                            "section": f"Slide {idx + 1}",
                            "text": "\n".join(slide_text)
                        })
            except Exception as e:
                pages_content.append({"page": 1, "section": "Document", "text": f"PPTX parse error: {e}"})

        else: # TXT fallback
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
                    pages_content.append({
                        "page": 1,
                        "section": "General Content",
                        "text": text
                    })
            except Exception as e:
                pages_content.append({"page": 1, "section": "Document", "text": f"TXT parse error: {e}"})

        return pages_content

document_processor = DocumentProcessor()
